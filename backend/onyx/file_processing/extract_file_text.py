import io
import json
import os
import re
import zipfile
from collections.abc import Callable
from collections.abc import Iterator
from collections.abc import Sequence
from email.parser import Parser as EmailParser
from enum import auto
from enum import IntFlag
from io import BytesIO
from pathlib import Path
from typing import Any
from typing import IO
from typing import NamedTuple
from zipfile import BadZipFile

import chardet
import docx  # type: ignore
import openpyxl  # type: ignore
import pptx  # type: ignore
from docx import Document as DocxDocument
from fastapi import UploadFile
from PIL import Image
from pypdf import PdfReader
from pypdf.errors import PdfStreamError

try:
    import frontmatter
    FRONTMATTER_AVAILABLE = True
except ImportError:
    FRONTMATTER_AVAILABLE = False

from onyx.configs.constants import FileOrigin
from onyx.configs.constants import ONYX_METADATA_FILENAME
from onyx.configs.llm_configs import get_image_extraction_and_analysis_enabled
from onyx.file_processing.html_utils import parse_html_page_basic
from onyx.file_processing.unstructured import get_unstructured_api_key
from onyx.file_processing.unstructured import unstructured_to_text
from onyx.file_store.file_store import FileStore
from onyx.utils.logger import setup_logger

logger = setup_logger()


class ExtractionResult(NamedTuple):
    """Structured result from text and image extraction from various file types."""

    text_content: str
    embedded_images: Sequence[tuple[bytes, str]]
    metadata: dict[str, Any]


# NOTE(rkuo): Unify this with upload_files_for_chat and file_valiation.py
TEXT_SECTION_SEPARATOR = "\n\n"

ACCEPTED_PLAIN_TEXT_FILE_EXTENSIONS = [
    ".txt",
    ".md",
    ".mdx",
    ".conf",
    ".log",
    ".json",
    ".csv",
    ".tsv",
    ".xml",
    ".yml",
    ".yaml",
]

ACCEPTED_DOCUMENT_FILE_EXTENSIONS = [
    ".pdf",
    ".docx",
    ".pptx",
    ".xlsx",
    ".eml",
    ".epub",
    ".html",
]

ACCEPTED_IMAGE_FILE_EXTENSIONS = [
    ".png",
    ".jpg",
    ".jpeg",
    ".webp",
]

ALL_ACCEPTED_FILE_EXTENSIONS = (
    ACCEPTED_PLAIN_TEXT_FILE_EXTENSIONS
    + ACCEPTED_DOCUMENT_FILE_EXTENSIONS
    + ACCEPTED_IMAGE_FILE_EXTENSIONS
)

IMAGE_MEDIA_TYPES = [
    "image/png",
    "image/jpeg",
    "image/webp",
]


class OnyxExtensionType(IntFlag):
    Plain = auto()
    Document = auto()
    Multimedia = auto()
    All = Plain | Document | Multimedia


def is_text_file_extension(file_name: str) -> bool:
    return any(file_name.endswith(ext) for ext in ACCEPTED_PLAIN_TEXT_FILE_EXTENSIONS)


def get_file_ext(file_path_or_name: str | Path) -> str:
    _, extension = os.path.splitext(file_path_or_name)
    return extension.lower()


def is_valid_media_type(media_type: str) -> bool:
    return media_type in IMAGE_MEDIA_TYPES


def is_accepted_file_ext(ext: str, ext_type: OnyxExtensionType) -> bool:
    if ext_type & OnyxExtensionType.Plain:
        if ext in ACCEPTED_PLAIN_TEXT_FILE_EXTENSIONS:
            return True

    if ext_type & OnyxExtensionType.Document:
        if ext in ACCEPTED_DOCUMENT_FILE_EXTENSIONS:
            return True

    if ext_type & OnyxExtensionType.Multimedia:
        if ext in ACCEPTED_IMAGE_FILE_EXTENSIONS:
            return True

    return False


def is_text_file(file: IO[bytes]) -> bool:
    """
    checks if the first 1024 bytes only contain printable or whitespace characters
    if it does, then we say it's a plaintext file
    """
    raw_data = file.read(1024)
    file.seek(0)
    text_chars = bytearray({7, 8, 9, 10, 12, 13, 27} | set(range(0x20, 0x100)) - {0x7F})
    return all(c in text_chars for c in raw_data)


def detect_encoding(file: IO[bytes]) -> str:
    raw_data = file.read(50000)
    file.seek(0)
    encoding = chardet.detect(raw_data)["encoding"] or "utf-8"
    return encoding


def is_macos_resource_fork_file(file_name: str) -> bool:
    return os.path.basename(file_name).startswith("._") and file_name.startswith(
        "__MACOSX"
    )


def load_files_from_zip(
    zip_file_io: IO,
    ignore_macos_resource_fork_files: bool = True,
    ignore_dirs: bool = True,
) -> Iterator[tuple[zipfile.ZipInfo, IO[Any]]]:
    """
    Iterates through files in a zip archive, yielding (ZipInfo, file handle) pairs.
    """
    with zipfile.ZipFile(zip_file_io, "r") as zip_file:
        for file_info in zip_file.infolist():
            if ignore_dirs and file_info.is_dir():
                continue

            if (
                ignore_macos_resource_fork_files
                and is_macos_resource_fork_file(file_info.filename)
            ) or file_info.filename == ONYX_METADATA_FILENAME:
                continue

            with zip_file.open(file_info.filename, "r") as subfile:
                # Try to match by exact filename first
                yield file_info, subfile


def _extract_onyx_metadata(line: str) -> dict | None:
    """
    Example: first line has:
        <!-- ONYX_METADATA={"title": "..."} -->
      or
        #ONYX_METADATA={"title":"..."}
    """
    html_comment_pattern = r"<!--\s*ONYX_METADATA=\{(.*?)\}\s*-->"
    hashtag_pattern = r"#ONYX_METADATA=\{(.*?)\}"

    html_comment_match = re.search(html_comment_pattern, line)
    hashtag_match = re.search(hashtag_pattern, line)

    if html_comment_match:
        json_str = html_comment_match.group(1)
    elif hashtag_match:
        json_str = hashtag_match.group(1)
    else:
        return None

    try:
        return json.loads("{" + json_str + "}")
    except json.JSONDecodeError:
        return None


def read_text_file(
    file: IO,
    encoding: str = "utf-8",
    errors: str = "replace",
    ignore_onyx_metadata: bool = True,
) -> tuple[str, dict]:
    """
    For plain text files. Optionally extracts Onyx metadata from the first line.
    """
    metadata = {}
    file_content_raw = ""
    for ind, line in enumerate(file):
        # decode
        try:
            line = line.decode(encoding) if isinstance(line, bytes) else line
        except UnicodeDecodeError:
            line = (
                line.decode(encoding, errors=errors)
                if isinstance(line, bytes)
                else line
            )

        # optionally parse metadata in the first line
        if ind == 0 and not ignore_onyx_metadata:
            potential_meta = _extract_onyx_metadata(line)
            if potential_meta is not None:
                metadata = potential_meta
                continue

        file_content_raw += line

    return file_content_raw, metadata


def pdf_to_text(file: IO[Any], pdf_pass: str | None = None) -> str:
    """
    Extract text from a PDF. For embedded images, a more complex approach is needed.
    This is a minimal approach returning text only.
    """
    text, _, _ = read_pdf_file(file, pdf_pass)
    return text


def read_pdf_file(
    file: IO[Any], pdf_pass: str | None = None, extract_images: bool = False
) -> tuple[str, dict[str, Any], Sequence[tuple[bytes, str]]]:
    """
    Returns the text, basic PDF metadata, and optionally extracted images.
    """
    metadata: dict[str, Any] = {}
    extracted_images: list[tuple[bytes, str]] = []
    try:
        pdf_reader = PdfReader(file)

        if pdf_reader.is_encrypted and pdf_pass is not None:
            decrypt_success = False
            try:
                decrypt_success = pdf_reader.decrypt(pdf_pass) != 0
            except Exception:
                logger.error("Unable to decrypt pdf")

            if not decrypt_success:
                return "", metadata, []
        elif pdf_reader.is_encrypted:
            logger.warning("No Password for an encrypted PDF, returning empty text.")
            return "", metadata, []

        # Basic PDF metadata
        if pdf_reader.metadata is not None:
            for key, value in pdf_reader.metadata.items():
                clean_key = key.lstrip("/")
                if isinstance(value, str) and value.strip():
                    metadata[clean_key] = value
                elif isinstance(value, list) and all(
                    isinstance(item, str) for item in value
                ):
                    metadata[clean_key] = ", ".join(value)

        text = TEXT_SECTION_SEPARATOR.join(
            page.extract_text() for page in pdf_reader.pages
        )

        if extract_images:
            for page_num, page in enumerate(pdf_reader.pages):
                for image_file_object in page.images:
                    image = Image.open(io.BytesIO(image_file_object.data))
                    img_byte_arr = io.BytesIO()
                    image.save(img_byte_arr, format=image.format)
                    img_bytes = img_byte_arr.getvalue()

                    image_name = (
                        f"page_{page_num + 1}_image_{image_file_object.name}."
                        f"{image.format.lower() if image.format else 'png'}"
                    )
                    extracted_images.append((img_bytes, image_name))

        return text, metadata, extracted_images

    except PdfStreamError:
        logger.exception("Invalid PDF file")
    except Exception:
        logger.exception("Failed to read PDF")

    return "", metadata, []


def docx_to_text_and_images(
    file: IO[Any], file_name: str = ""
) -> tuple[str, Sequence[tuple[bytes, str]]]:
    """
    Extract text from a docx. If embed_images=True, also extract inline images.
    Return (text_content, list_of_images).
    """
    paragraphs = []
    embedded_images: list[tuple[bytes, str]] = []

    try:
        doc = docx.Document(file)
    except BadZipFile as e:
        logger.warning(f"Failed to extract text from {file_name or 'docx file'}: {e}")
        return "", []

    # Grab text from paragraphs
    for paragraph in doc.paragraphs:
        paragraphs.append(paragraph.text)

    # Reset position so we can re-load the doc (python-docx has read the stream)
    # Note: if python-docx has fully consumed the stream, you may need to open it again from memory.
    # For large docs, a more robust approach is needed.
    # This is a simplified example.

    for rel_id, rel in doc.part.rels.items():
        if "image" in rel.reltype:
            # image is typically in rel.target_part.blob
            image_bytes = rel.target_part.blob
            image_name = rel.target_part.partname
            # store
            embedded_images.append((image_bytes, os.path.basename(str(image_name))))

    text_content = "\n".join(paragraphs)
    return text_content, embedded_images


def pptx_to_text(file: IO[Any], file_name: str = "") -> str:
    try:
        presentation = pptx.Presentation(file)
    except BadZipFile as e:
        error_str = f"Failed to extract text from {file_name or 'pptx file'}: {e}"
        logger.warning(error_str)
        return ""
    text_content = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        slide_text = f"\nSlide {slide_number}:\n"
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text += shape.text + "\n"
        text_content.append(slide_text)
    return TEXT_SECTION_SEPARATOR.join(text_content)


def xlsx_to_text(file: IO[Any], file_name: str = "") -> str:
    try:
        workbook = openpyxl.load_workbook(file, read_only=True)
    except BadZipFile as e:
        error_str = f"Failed to extract text from {file_name or 'xlsx file'}: {e}"
        if file_name.startswith("~"):
            logger.debug(error_str + " (this is expected for files with ~)")
        else:
            logger.warning(error_str)
        return ""
    except Exception as e:
        if "File contains no valid workbook part" in str(e):
            logger.error(
                f"Failed to extract text from {file_name or 'xlsx file'}. This happens due to a bug in openpyxl. {e}"
            )
            return ""
        raise e

    text_content = []
    for sheet in workbook.worksheets:
        rows = []
        num_empty_consecutive_rows = 0
        for row in sheet.iter_rows(min_row=1, values_only=True):
            row_str = ",".join(str(cell or "") for cell in row)

            # Only add the row if there are any values in the cells
            if len(row_str) >= len(row):
                rows.append(row_str)
                num_empty_consecutive_rows = 0
            else:
                num_empty_consecutive_rows += 1

            if num_empty_consecutive_rows > 100:
                # handle massive excel sheets with mostly empty cells
                logger.warning(
                    f"Found {num_empty_consecutive_rows} empty rows in {file_name},"
                    " skipping rest of file"
                )
                break
        sheet_str = "\n".join(rows)
        text_content.append(sheet_str)
    return TEXT_SECTION_SEPARATOR.join(text_content)


def eml_to_text(file: IO[Any]) -> str:
    encoding = detect_encoding(file)
    text_file = io.TextIOWrapper(file, encoding=encoding)
    parser = EmailParser()
    message = parser.parse(text_file)

    text_content = []
    for part in message.walk():
        if part.get_content_type().startswith("text/plain"):
            payload = part.get_payload()
            if isinstance(payload, str):
                text_content.append(payload)
            elif isinstance(payload, list):
                text_content.extend(item for item in payload if isinstance(item, str))
            else:
                logger.warning(f"Unexpected payload type: {type(payload)}")
    return TEXT_SECTION_SEPARATOR.join(text_content)


def epub_to_text(file: IO[Any]) -> str:
    with zipfile.ZipFile(file) as epub:
        text_content = []
        for item in epub.infolist():
            if item.filename.endswith(".xhtml") or item.filename.endswith(".html"):
                with epub.open(item) as html_file:
                    text_content.append(parse_html_page_basic(html_file))
        return TEXT_SECTION_SEPARATOR.join(text_content)


def file_io_to_text(file: IO[Any]) -> str:
    encoding = detect_encoding(file)
    file_content, _ = read_text_file(file, encoding=encoding)
    return file_content


def markdown_to_text_and_metadata(
    file: IO[Any], 
    file_name: str = ""
) -> ExtractionResult:
    """
    专门的Markdown解析器，支持YAML front matter和链接提取
    
    功能：
    - 解析YAML front matter metadata
    - 提取内联Markdown链接  
    - 保持文档结构用于分块
    - 映射metadata到Document模型
    
    Args:
        file: 文件IO对象
        file_name: 文件名（用于错误处理）
        
    Returns:
        ExtractionResult包含text_content, embedded_images, metadata
    """
    
    if not FRONTMATTER_AVAILABLE:
        logger.warning(f"python-frontmatter not available, falling back to plain text processing for {file_name}")
        file.seek(0)
        encoding = detect_encoding(file)
        text_content_raw, file_metadata = read_text_file(
            file, encoding=encoding, ignore_onyx_metadata=False
        )
        return ExtractionResult(
            text_content=text_content_raw,
            embedded_images=[],
            metadata=file_metadata
        )
    
    try:
        # 读取并解析front matter
        file.seek(0)
        content = file.read()
        if isinstance(content, bytes):
            content = content.decode('utf-8', errors='replace')
        
        post = frontmatter.loads(content)
        
        # 提取front matter metadata和markdown内容
        front_matter_metadata = post.metadata
        markdown_content = post.content
        
        # 提取内联链接 [text](url)
        link_pattern = r'\[([^\]]+)\]\(([^)]+)\)'
        inline_links = re.findall(link_pattern, markdown_content)
        
        # 处理metadata
        processed_metadata = {}
        
        # 标准字段处理
        if 'title' in front_matter_metadata:
            processed_metadata['title'] = str(front_matter_metadata['title'])
        
        # 处理作者信息（支持字符串和对象格式）
        if 'author' in front_matter_metadata:
            author = front_matter_metadata['author']
            if isinstance(author, dict):
                processed_metadata['author_name'] = author.get('name', '')
                if author.get('email'):
                    processed_metadata['author_email'] = author.get('email')
                if author.get('link'):
                    processed_metadata['author_link'] = author.get('link')
            else:
                processed_metadata['author_name'] = str(author)
        
        # 处理日期
        if 'date' in front_matter_metadata:
            date_value = front_matter_metadata['date']
            if hasattr(date_value, 'strftime'):  # datetime对象
                processed_metadata['date'] = date_value.strftime('%Y-%m-%d')
            else:
                processed_metadata['date'] = str(date_value)
        
        # 处理标签
        if 'tags' in front_matter_metadata:
            tags = front_matter_metadata['tags']
            if isinstance(tags, list):
                processed_metadata['tags'] = [str(tag) for tag in tags]
            elif isinstance(tags, str):
                # 支持逗号分隔的标签字符串
                processed_metadata['tags'] = [tag.strip() for tag in tags.split(',')]
            else:
                processed_metadata['tags'] = [str(tags)]
        
        # 关键：source_link从front matter的link字段获取
        if 'link' in front_matter_metadata:
            processed_metadata['source_link'] = str(front_matter_metadata['link'])
        
        # 处理其他标准字段
        for field in ['category', 'description', 'version', 'status']:
            if field in front_matter_metadata:
                processed_metadata[field] = str(front_matter_metadata[field])
        
        # 处理贡献者
        if 'contributors' in front_matter_metadata:
            contributors = front_matter_metadata['contributors']
            if isinstance(contributors, list):
                processed_metadata['contributors'] = [str(c) for c in contributors]
            else:
                processed_metadata['contributors'] = [str(contributors)]
        
        # 添加内联链接作为metadata
        if inline_links:
            processed_metadata['inline_links'] = [
                f"{text}|{url}" for text, url in inline_links
            ]
        
        # 添加其他自定义字段
        standard_fields = {
            'title', 'author', 'date', 'tags', 'link', 'category', 
            'description', 'version', 'status', 'contributors'
        }
        for key, value in front_matter_metadata.items():
            if key not in standard_fields:
                processed_metadata[f"custom_{key}"] = str(value)
        
        return ExtractionResult(
            text_content=markdown_content,
            embedded_images=[],  # Markdown文件不包含二进制图像
            metadata=processed_metadata
        )
        
    except Exception as e:
        logger.exception(f"Failed to parse Markdown metadata from {file_name}: {e}")
        # 降级到普通文本处理
        file.seek(0)
        encoding = detect_encoding(file)
        text_content_raw, file_metadata = read_text_file(
            file, encoding=encoding, ignore_onyx_metadata=False
        )
        return ExtractionResult(
            text_content=text_content_raw,
            embedded_images=[],
            metadata=file_metadata
        )


def extract_file_text(
    file: IO[Any],
    file_name: str,
    break_on_unprocessable: bool = True,
    extension: str | None = None,
) -> str:
    """
    Legacy function that returns *only text*, ignoring embedded images.
    For backward-compatibility in code that only wants text.

    NOTE: Ignoring seems to be defined as returning an empty string for files it can't
    handle (such as images).
    """
    extension_to_function: dict[str, Callable[[IO[Any]], str]] = {
        ".pdf": pdf_to_text,
        ".docx": lambda f: docx_to_text_and_images(f)[0],  # no images
        ".pptx": pptx_to_text,
        ".xlsx": xlsx_to_text,
        ".eml": eml_to_text,
        ".epub": epub_to_text,
        ".html": parse_html_page_basic,
    }

    try:
        if get_unstructured_api_key():
            try:
                return unstructured_to_text(file, file_name)
            except Exception as unstructured_error:
                logger.error(
                    f"Failed to process with Unstructured: {str(unstructured_error)}. "
                    "Falling back to normal processing."
                )
        if extension is None:
            extension = get_file_ext(file_name)

        if is_accepted_file_ext(
            extension, OnyxExtensionType.Plain | OnyxExtensionType.Document
        ):
            func = extension_to_function.get(extension, file_io_to_text)
            file.seek(0)
            return func(file)

        # If unknown extension, maybe it's a text file
        file.seek(0)
        if is_text_file(file):
            return file_io_to_text(file)

        raise ValueError("Unknown file extension or not recognized as text data")

    except Exception as e:
        if break_on_unprocessable:
            raise RuntimeError(
                f"Failed to process file {file_name or 'Unknown'}: {str(e)}"
            ) from e
        logger.warning(f"Failed to process file {file_name or 'Unknown'}: {str(e)}")
        return ""


def extract_text_and_images(
    file: IO[Any],
    file_name: str,
    pdf_pass: str | None = None,
) -> ExtractionResult:
    """
    Primary new function for the updated connector.
    Returns structured extraction result with text content, embedded images, and metadata.
    """

    try:
        # Attempt unstructured if env var is set
        if get_unstructured_api_key():
            # If the user doesn't want embedded images, unstructured is fine
            file.seek(0)
            text_content = unstructured_to_text(file, file_name)
            return ExtractionResult(
                text_content=text_content, embedded_images=[], metadata={}
            )

        extension = get_file_ext(file_name)

        # docx example for embedded images
        if extension == ".docx":
            file.seek(0)
            text_content, images = docx_to_text_and_images(file)
            return ExtractionResult(
                text_content=text_content, embedded_images=images, metadata={}
            )

        # PDF example: we do not show complicated PDF image extraction here
        # so we simply extract text for now and skip images.
        if extension == ".pdf":
            file.seek(0)
            text_content, pdf_metadata, images = read_pdf_file(
                file,
                pdf_pass,
                extract_images=get_image_extraction_and_analysis_enabled(),
            )
            return ExtractionResult(
                text_content=text_content, embedded_images=images, metadata=pdf_metadata
            )

        # For PPTX, XLSX, EML, etc., we do not show embedded image logic here.
        # You can do something similar to docx if needed.
        if extension == ".pptx":
            file.seek(0)
            return ExtractionResult(
                text_content=pptx_to_text(file, file_name=file_name),
                embedded_images=[],
                metadata={},
            )

        if extension == ".xlsx":
            file.seek(0)
            return ExtractionResult(
                text_content=xlsx_to_text(file, file_name=file_name),
                embedded_images=[],
                metadata={},
            )

        if extension == ".eml":
            file.seek(0)
            return ExtractionResult(
                text_content=eml_to_text(file), embedded_images=[], metadata={}
            )

        if extension == ".epub":
            file.seek(0)
            return ExtractionResult(
                text_content=epub_to_text(file), embedded_images=[], metadata={}
            )

        if extension == ".html":
            file.seek(0)
            return ExtractionResult(
                text_content=parse_html_page_basic(file),
                embedded_images=[],
                metadata={},
            )

        # Special handling for Markdown files with frontmatter support
        if extension in [".md", ".mdx"]:
            file.seek(0)
            return markdown_to_text_and_metadata(file, file_name)

        # If we reach here and it's a recognized text extension
        if is_text_file_extension(file_name):
            file.seek(0)
            encoding = detect_encoding(file)
            text_content_raw, file_metadata = read_text_file(
                file, encoding=encoding, ignore_onyx_metadata=False
            )
            return ExtractionResult(
                text_content=text_content_raw,
                embedded_images=[],
                metadata=file_metadata,
            )

        # If it's an image file or something else, we do not parse embedded images from them
        # just return empty text
        file.seek(0)
        return ExtractionResult(text_content="", embedded_images=[], metadata={})

    except Exception as e:
        logger.exception(f"Failed to extract text/images from {file_name}: {e}")
        return ExtractionResult(text_content="", embedded_images=[], metadata={})


def convert_docx_to_txt(file: UploadFile, file_store: FileStore) -> str:
    """
    Helper to convert docx to a .txt file in the same filestore.
    """
    file.file.seek(0)
    docx_content = file.file.read()
    doc = DocxDocument(BytesIO(docx_content))

    # Extract text from the document
    all_paras = [p.text for p in doc.paragraphs]
    text_content = "\n".join(all_paras)

    file_id = file_store.save_file(
        content=BytesIO(text_content.encode("utf-8")),
        display_name=file.filename,
        file_origin=FileOrigin.CONNECTOR,
        file_type="text/plain",
    )
    return file_id


def docx_to_txt_filename(file_path: str) -> str:
    return file_path.rsplit(".", 1)[0] + ".txt"
